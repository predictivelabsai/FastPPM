"""Export a report (ordered blocks) to PDF, DOCX or PPTX.

Blocks are {type, content} where content is HTML. Inline HTML is reduced to text
(stdlib html.parser, no regex); bullets are split into items. PDF via fpdf2,
DOCX via python-docx, PPTX via python-pptx — all pure-Python (no browser), so
this runs in the container.
"""

from __future__ import annotations

import io
from html.parser import HTMLParser

NAVY = (107, 23, 102)   # #123B5D
SLATE = (72, 72, 79)


# ── HTML → text helpers (no regex) ──────────────────────────────────────────

class _Text(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts = []

    def handle_data(self, d):
        self.parts.append(d)


def html_to_text(html: str) -> str:
    p = _Text()
    p.feed(html or "")
    return " ".join("".join(p.parts).split())


class _Items(HTMLParser):
    def __init__(self):
        super().__init__()
        self.items = []
        self._cur = None

    def handle_starttag(self, tag, attrs):
        if tag == "li":
            self._cur = []

    def handle_endtag(self, tag):
        if tag == "li" and self._cur is not None:
            self.items.append(" ".join("".join(self._cur).split()))
            self._cur = None

    def handle_data(self, d):
        if self._cur is not None:
            self._cur.append(d)


def list_items(html: str) -> list[str]:
    p = _Items()
    p.feed(html or "")
    return [i for i in p.items if i]


def _latin1(s: str) -> str:
    """Map common unicode punctuation to latin-1 so fpdf2's core fonts don't choke."""
    repl = {"—": "-", "–": "-", "•": "-", "“": '"', "”": '"', "‘": "'", "’": "'",
            "…": "...", "×": "x", "→": "->", "·": "-", "€": "EUR", " ": " "}
    for k, v in repl.items():
        s = s.replace(k, v)
    return s.encode("latin-1", "ignore").decode("latin-1")


# ── PDF (fpdf2) ─────────────────────────────────────────────────────────────

def pdf_bytes(title: str, blocks: list[dict]) -> bytes:
    from fpdf import FPDF
    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_auto_page_break(True, margin=18)
    pdf.add_page()
    pdf.set_margins(20, 18, 20)

    for b in blocks:
        t, c = b["type"], b.get("content", "")
        if t == "divider":
            pdf.ln(2); y = pdf.get_y()
            pdf.set_draw_color(220, 217, 228); pdf.line(20, y, 190, y); pdf.ln(4)
            continue
        text = _latin1(html_to_text(c))
        if t == "heading1":
            pdf.set_font("Helvetica", "B", 21); pdf.set_text_color(*NAVY)
            pdf.multi_cell(0, 9, text); pdf.ln(2)
        elif t == "heading2":
            pdf.ln(2); pdf.set_font("Helvetica", "B", 15); pdf.set_text_color(*NAVY)
            pdf.multi_cell(0, 7, text); pdf.ln(1)
        elif t == "heading3":
            pdf.set_font("Helvetica", "B", 12); pdf.set_text_color(*SLATE)
            pdf.multi_cell(0, 6, text); pdf.ln(1)
        elif t == "bullet":
            pdf.set_font("Helvetica", "", 11); pdf.set_text_color(*SLATE)
            for item in list_items(c) or ([text] if text else []):
                pdf.set_x(22)
                pdf.multi_cell(0, 6, "-  " + _latin1(item))
            pdf.ln(1)
        else:  # paragraph
            if text:
                pdf.set_font("Helvetica", "", 11); pdf.set_text_color(*SLATE)
                pdf.multi_cell(0, 6, text); pdf.ln(1)
    out = pdf.output()
    return bytes(out)


# ── DOCX (python-docx) ──────────────────────────────────────────────────────

def docx_bytes(title: str, blocks: list[dict]) -> bytes:
    from docx import Document
    doc = Document()
    for b in blocks:
        t, c = b["type"], b.get("content", "")
        if t == "divider":
            doc.add_paragraph("")
            continue
        text = html_to_text(c)
        if t == "heading1":
            doc.add_heading(text, level=0)
        elif t == "heading2":
            doc.add_heading(text, level=1)
        elif t == "heading3":
            doc.add_heading(text, level=2)
        elif t == "bullet":
            for item in list_items(c) or ([text] if text else []):
                doc.add_paragraph(item, style="List Bullet")
        elif text:
            doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── PPTX (python-pptx) ──────────────────────────────────────────────────────

def pptx_bytes(title: str, blocks: list[dict]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_layout = prs.slide_layouts[5]   # title only
    blank = prs.slide_layouts[6]

    # Title slide.
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = html_to_text(blocks[0]["content"]) if blocks else (title or "Report")

    def _new_section(heading):
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
        p = tb.text_frame.paragraphs[0]
        p.text = heading
        p.font.size = Pt(30)
        p.font.bold = True
        body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.4))
        return body.text_frame

    tf = None
    for b in blocks[1:]:
        t, c = b["type"], b.get("content", "")
        text = html_to_text(c)
        if t in ("heading1", "heading2"):
            tf = _new_section(text)
            tf.word_wrap = True
            tf.paragraphs[0].text = ""
        elif tf is None:
            tf = _new_section("Overview")
        if t == "heading3":
            para = tf.add_paragraph(); para.text = text; para.font.bold = True; para.font.size = Pt(18)
        elif t == "bullet":
            for item in list_items(c):
                para = tf.add_paragraph(); para.text = "•  " + item; para.font.size = Pt(16)
        elif t == "paragraph" and text:
            para = tf.add_paragraph(); para.text = text; para.font.size = Pt(16)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


XLSX_MEDIA = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

MEDIA = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def export(fmt: str, title: str, blocks: list[dict]) -> tuple[bytes, str]:
    """Return (bytes, media_type) for the requested format."""
    if fmt == "docx":
        return docx_bytes(title, blocks), MEDIA["docx"]
    if fmt == "pptx":
        return pptx_bytes(title, blocks), MEDIA["pptx"]
    return pdf_bytes(title, blocks), MEDIA["pdf"]
