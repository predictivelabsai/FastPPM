"""Per-format text + table extraction for the ingestion engine.

Turns a raw file (PDF / XLSX / PPTX / DOCX) into a normalised payload:

    {"text": "<all flowing text>", "tables": [[[cell,...], ...], ...]}

Kept dependency-light and defensive — a malformed file yields whatever could be
read plus an empty table list, never an exception that aborts a bulk upload.
"""

from __future__ import annotations

import os

SUPPORTED = {".pdf", ".xlsx", ".xls", ".pptx", ".docx"}


def file_type(name: str) -> str:
    ext = os.path.splitext(name)[1].lower()
    return {".pdf": "pdf", ".xlsx": "xlsx", ".xls": "xlsx",
            ".pptx": "pptx", ".docx": "docx"}.get(ext, ext.lstrip(".") or "unknown")


def extract(path: str) -> dict:
    """Dispatch on extension. Returns {text, tables}."""
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pdf":
            return _pdf(path)
        if ext in (".xlsx", ".xls"):
            return _xlsx(path)
        if ext == ".pptx":
            return _pptx(path)
        if ext == ".docx":
            return _docx(path)
    except Exception as e:  # noqa: BLE001 — never abort a bulk upload
        return {"text": f"[extraction error: {e}]", "tables": []}
    return {"text": "", "tables": []}


def _pdf(path: str) -> dict:
    import pdfplumber
    texts, tables = [], []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text() or ""
            if t:
                texts.append(t)
            for tbl in page.extract_tables() or []:
                rows = [[(c or "").strip() for c in row] for row in tbl]
                if rows:
                    tables.append(rows)
    return {"text": "\n".join(texts), "tables": tables}


def _xlsx(path: str) -> dict:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    texts, tables = [], []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(cells):
                rows.append(cells)
        if rows:
            texts.append(f"# Sheet: {ws.title}\n" +
                         "\n".join("\t".join(r) for r in rows))
            tables.append(rows)
    wb.close()
    return {"text": "\n\n".join(texts), "tables": tables}


def _pptx(path: str) -> dict:
    from pptx import Presentation
    prs = Presentation(path)
    texts, tables = [], []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"# Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                rows = [[cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows]
                if rows:
                    tables.append(rows)
                    parts.append("\n".join("\t".join(r) for r in rows))
        texts.append("\n".join(parts))
    return {"text": "\n\n".join(texts), "tables": tables}


def _docx(path: str) -> dict:
    from docx import Document
    doc = Document(path)
    texts = [p.text for p in doc.paragraphs if p.text.strip()]
    tables = []
    for tbl in doc.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
        if rows:
            tables.append(rows)
            texts.append("\n".join("\t".join(r) for r in rows))
    return {"text": "\n".join(texts), "tables": tables}
