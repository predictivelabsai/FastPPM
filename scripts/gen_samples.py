"""Generate messy sample source documents for the ingestion demo.

Writes a PDF status report, an XLSX tracker, a PPTX board deck and a DOCX charter
into ``samples/`` — each about a transformation initiative NOT already in the
seeded master view, with deliberately varying column names and date formats so
the normalisation engine has something to reconcile.
"""

from __future__ import annotations

from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
SAMPLES = ROOT / "samples"


def _pdf():
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 15)
    pdf.cell(0, 10, "Digital Customer Portal - Monthly Status Report", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 10)
    pdf.cell(0, 7, "Owner: R. Mensah   |   Workstream: Customer Experience   |   May 2026",
             new_x="LMARGIN", new_y="NEXT")
    pdf.ln(2)
    pdf.multi_cell(0, 6, "Overall status: At risk. The portal build is progressing but a "
                   "vendor integration slip threatens the July go-live. Target benefit is "
                   "GBP 2.4m of revenue uplift, of which GBP 0.6m realised to date.")
    pdf.ln(2)

    pdf.set_font("Helvetica", "B", 11)
    pdf.cell(0, 8, "Milestones", new_x="LMARGIN", new_y="NEXT")
    headers = ["Milestone", "Owner", "Due date", "Actual", "% Complete", "Status"]
    widths = [62, 26, 24, 24, 22, 22]
    rows = [
        ["Discovery & requirements", "R. Mensah", "15/02/2026", "18/02/2026", "100", "Done"],
        ["UX design sign-off", "P. Costa", "20/03/2026", "27/03/2026", "100", "Done"],
        ["Portal build (MVP)", "S. Yilmaz", "30/05/2026", "", "70", "In progress"],
        ["Vendor payment integration", "S. Yilmaz", "15/06/2026", "", "35", "At risk"],
        ["Go-live", "R. Mensah", "15/07/2026", "", "0", "Planned"],
    ]
    _pdf_table(pdf, headers, widths, rows)

    pdf.ln(3)
    pdf.set_font("Helvetica", "B", 11)
    pdf.cell(0, 8, "Risks", new_x="LMARGIN", new_y="NEXT")
    _pdf_table(pdf, ["Risk", "Probability", "Impact", "Mitigation"], [70, 26, 22, 62], [
        ["Vendor integration delay", "4", "5", "Escalate; dual-source provider"],
        ["Scope creep on MVP", "3", "3", "Strict change control"],
    ])
    out = SAMPLES / "Digital_Customer_Portal_StatusReport.pdf"
    pdf.output(str(out))
    return out


def _pdf_table(pdf, headers, widths, rows):
    pdf.set_font("Helvetica", "B", 8)
    for h, w in zip(headers, widths):
        pdf.cell(w, 7, h, border=1)
    pdf.ln()
    pdf.set_font("Helvetica", "", 8)
    for row in rows:
        for c, w in zip(row, widths):
            pdf.cell(w, 7, str(c)[:int(w / 1.7)], border=1)
        pdf.ln()


def _xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Milestones"
    ws.append(["Task", "Lead", "Planned", "Forecast", "Progress", "RAG"])
    for r in [
        ["Data pipeline build", "A. Schmidt", "2026-03-01", "2026-03-10", "100", "Green"],
        ["Model training & validation", "G. Petrova", "2026-04-15", "2026-04-20", "90", "Green"],
        ["Forecast accuracy >85%", "G. Petrova", "2026-05-30", "", "60", "Amber"],
        ["Integration with planning", "A. Schmidt", "2026-07-01", "", "20", "Amber"],
        ["Roll-out to 3 regions", "M. Haddad", "2026-09-15", "", "0", "Planned"],
    ]:
        ws.append(r)
    vs = wb.create_sheet("Value")
    vs.append(["Value driver", "Category", "Target", "Realised", "Quarter"])
    for r in [
        ["Inventory reduction", "Cost savings", "1,800,000", "700,000", "Q2 2026"],
        ["Stockout avoidance", "Revenue", "900,000", "250,000", "Q2 2026"],
    ]:
        vs.append(r)
    out = SAMPLES / "SupplyChain_AI_Forecasting_Tracker.xlsx"
    wb.save(out)
    return out


def _pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[5]

    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "Finance Automation - Board Update (AI use case)"

    s2 = prs.slides.add_slide(blank)
    tb = s2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tb.text_frame.text = ("Owner: D. Iqbal. Status: On track. Automating accounts payable "
                          "with an LLM document agent. Target EBITDA impact GBP 1.3m.")
    rows, cols = 3, 5
    tbl = s2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.6), Inches(9), Inches(2)).table
    for j, h in enumerate(["Value driver", "Category", "Target", "Captured", "Period"]):
        tbl.cell(0, j).text = h
    data = [["Invoice processing cost", "Cost savings", "900,000", "500,000", "Q2 2026"],
            ["Early-payment discounts", "EBITDA", "400,000", "150,000", "Q2 2026"]]
    for i, row in enumerate(data, 1):
        for j, v in enumerate(row):
            tbl.cell(i, j).text = v
    out = SAMPLES / "Finance_Automation_Board_Update.pptx"
    prs.save(out)
    return out


def _docx():
    from docx import Document
    doc = Document()
    doc.add_heading("HR Shared Services - Programme Charter", level=1)
    doc.add_paragraph("Owner: N. Abara. Workstream: Operations. Status: In progress. "
                      "Consolidating HR operations into a shared-services centre. "
                      "Target cost savings GBP 2.0m.")
    doc.add_heading("Milestones", level=2)
    t = doc.add_table(rows=1, cols=5)
    for j, h in enumerate(["Deliverable", "Owner", "Target date", "Actual date", "Status"]):
        t.rows[0].cells[j].text = h
    for r in [["Operating model design", "N. Abara", "28 Feb 2026", "05 Mar 2026", "Done"],
              ["Process standardisation", "L. Bianchi", "30 Apr 2026", "", "In progress"],
              ["Systems migration", "A. Schmidt", "31 Jul 2026", "", "Planned"]]:
        cells = t.add_row().cells
        for j, v in enumerate(r):
            cells[j].text = v
    doc.add_heading("Risks", level=2)
    rt = doc.add_table(rows=1, cols=4)
    for j, h in enumerate(["Risk", "Probability", "Impact", "Mitigation"]):
        rt.rows[0].cells[j].text = h
    for r in [["Change resistance from BUs", "3", "4", "Stakeholder engagement plan"],
              ["Key-person dependency", "3", "3", "Cross-training"]]:
        cells = rt.add_row().cells
        for j, v in enumerate(r):
            cells[j].text = v
    out = SAMPLES / "HR_SharedServices_Charter.docx"
    doc.save(out)
    return out


def generate() -> list[Path]:
    SAMPLES.mkdir(exist_ok=True)
    return [_pdf(), _xlsx(), _pptx(), _docx()]


if __name__ == "__main__":
    for p in generate():
        print("wrote", p.relative_to(ROOT))
