"""Generate a few *new* synthetic source documents into data/incoming/.

These cover initiatives NOT in the seeded programme, so ingesting them through
the UI demonstrates the engine discovering and adding new initiatives to the
master repository (vs. the seeded samples/ docs). Deliberately messy: varying
column names, date formats and structures.
"""

from __future__ import annotations

from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "data" / "incoming"


def _pdf_table(pdf, headers, widths, rows):
    pdf.set_font("Helvetica", "B", 8)
    for h, w in zip(headers, widths):
        pdf.cell(w, 7, h, border=1)
    pdf.ln()
    pdf.set_font("Helvetica", "", 8)
    for row in rows:
        for c, w in zip(row, widths):
            pdf.cell(w, 7, str(c)[:int(w / 1.7)], border=1)
        pdf.ln()


def _pdf():
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 15)
    pdf.cell(0, 10, "Cybersecurity & Resilience Uplift - Status Report",
             new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 10)
    pdf.cell(0, 7, "Owner: T. Novak   |   Workstream: Technology   |   June 2026",
             new_x="LMARGIN", new_y="NEXT")
    pdf.ln(2)
    pdf.multi_cell(0, 6, "Overall status: At risk. Hardening the estate after the "
                   "diligence findings. Target benefit GBP 1.1m of avoided loss / "
                   "cost savings, GBP 0.2m realised so far.")
    pdf.ln(2)
    pdf.set_font("Helvetica", "B", 11)
    pdf.cell(0, 8, "Milestones", new_x="LMARGIN", new_y="NEXT")
    _pdf_table(pdf, ["Milestone", "Owner", "Due date", "Actual", "% Complete", "Status"],
               [62, 26, 24, 24, 22, 22], [
        ["Threat assessment & baseline", "T. Novak", "10/02/2026", "14/02/2026", "100", "Done"],
        ["MFA & identity rollout", "A. Schmidt", "31/03/2026", "05/04/2026", "100", "Done"],
        ["Endpoint hardening", "T. Novak", "31/05/2026", "", "55", "In progress"],
        ["SOC / monitoring stand-up", "M. Haddad", "30/06/2026", "", "30", "At risk"],
        ["Pen-test & remediation", "T. Novak", "31/08/2026", "", "0", "Planned"],
    ])
    pdf.ln(3)
    pdf.set_font("Helvetica", "B", 11)
    pdf.cell(0, 8, "Risks", new_x="LMARGIN", new_y="NEXT")
    _pdf_table(pdf, ["Risk", "Probability", "Impact", "Mitigation"], [70, 26, 22, 62], [
        ["Legacy systems cannot be patched", "4", "5", "Compensating controls; isolate"],
        ["Security skills shortage", "3", "4", "MSSP partner for SOC"],
    ])
    OUT.mkdir(parents=True, exist_ok=True)
    out = OUT / "Cybersecurity_Resilience_Uplift_StatusReport.pdf"
    pdf.output(str(out))
    return out


def _xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Milestones"
    ws.append(["Task", "Lead", "Planned", "Forecast", "Progress", "RAG"])
    for r in [
        ["Discovery & data audit", "L. Bianchi", "2026-02-15", "2026-02-20", "100", "Green"],
        ["Elasticity model rebuild", "G. Petrova", "2026-04-30", "2026-05-08", "85", "Amber"],
        ["Pilot on 2 categories", "L. Bianchi", "2026-06-30", "", "40", "Amber"],
        ["Full price-book rollout", "J. Okafor", "2026-09-30", "", "0", "Planned"],
    ]:
        ws.append(r)
    vs = wb.create_sheet("Value")
    vs.append(["Value driver", "Category", "Target", "Realised", "Quarter"])
    for r in [
        ["Margin uplift from repricing", "EBITDA", "1,600,000", "350,000", "Q2 2026"],
        ["Discount leakage reduction", "Cost savings", "600,000", "180,000", "Q2 2026"],
    ]:
        vs.append(r)
    OUT.mkdir(parents=True, exist_ok=True)
    out = OUT / "Pricing_Analytics_Refresh_Tracker.xlsx"
    wb.save(out)
    return out


def _pptx():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[5]
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "Manufacturing Quality AI - Board Update (AI use case)"
    s2 = prs.slides.add_slide(blank)
    tb = s2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tb.text_frame.text = ("Owner: N. Abara. Status: On track. Computer-vision defect "
                          "detection on Line 2. Target cost savings GBP 1.4m.")
    rows, cols = 3, 5
    tbl = s2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.6), Inches(9), Inches(2)).table
    for j, h in enumerate(["Value driver", "Category", "Target", "Captured", "Period"]):
        tbl.cell(0, j).text = h
    for i, row in enumerate([["Scrap reduction", "Cost savings", "900,000", "300,000", "Q2 2026"],
                             ["Rework avoidance", "EBITDA", "500,000", "120,000", "Q2 2026"]], 1):
        for j, v in enumerate(row):
            tbl.cell(i, j).text = v
    OUT.mkdir(parents=True, exist_ok=True)
    out = OUT / "Manufacturing_QualityAI_BoardUpdate.pptx"
    prs.save(out)
    return out


def generate():
    return [_pdf(), _xlsx(), _pptx()]


if __name__ == "__main__":
    for p in generate():
        print("wrote", p.relative_to(ROOT))
