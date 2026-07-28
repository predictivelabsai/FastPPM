"""Build the FastPPM user guide as a branded **landscape 16:9 slide deck**.

HTML/CSS controls the whole layout (FastPPM palette). Each `.slide` is rendered in
headless Chromium → PNG, then stitched into:

  - docs/fastppm_user_guide_YYYY-MM-DD.pdf   (landscape; fpdf2)
  - docs/fastppm_user_guide_YYYY-MM-DD.pptx  (16:9; python-pptx)
  - docs/fastppm_user_guide_YYYY-MM-DD.html  (the source deck)

Run:  python -m scripts.generate_user_guide
"""

from __future__ import annotations

import base64
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
SHOTS = DOCS / "screenshots"
W, H = 1280, 720  # 16:9

# ── Slide content ───────────────────────────────────────────────────────────

def _bullets(items):
    return "<ul>" + "".join(f"<li>{x}</li>" for x in items) + "</ul>"


def _cards(cards):
    return "<div class='cards'>" + "".join(
        f"<div class='card'><div class='ct'>{t}</div><div class='cb'>{b}</div></div>"
        for t, b in cards) + "</div>"


def _journey(steps):
    """A numbered happy-path: step badge + title + one-line description, laid out
    as a 3×2 grid of connected cards."""
    out = []
    for i, (title, body) in enumerate(steps, 1):
        out.append(
            f"<div class='step'><div class='sn'>{i}</div>"
            f"<div class='st'>{title}</div><div class='sb'>{body}</div></div>")
    return "<div class='journey'>" + "".join(out) + "</div>"


def _img_uri(name: str) -> str:
    p = SHOTS / name
    if not p.exists():
        return ""
    return "data:image/png;base64," + base64.b64encode(p.read_bytes()).decode()


def _shot(name: str, caption: str) -> str:
    """A slide body that is a screenshot + a one-line caption."""
    uri = _img_uri(name)
    if not uri:
        return f"<div class='cap'>{caption}</div><p>(screenshot {name} missing)</p>"
    return (f"<div class='shot'><img src='{uri}'></div>"
            f"<div class='cap'>{caption}</div>")


SLIDES = [
    ("title", """
      <div class='hero'>
        <div class='brand'>Value&nbsp;Hub</div>
        <div class='sub'>Value Creation Plan · Transformation Office</div>
        <div class='tag'>A conversational-first tool to plan, track and report every
          transformation initiative and internal AI use case — for a PE portfolio company.</div>
        <div class='foot'>User Guide</div>
      </div>"""),
    ("Your typical journey", "<div class='lead'>From a messy status report to a "
        "board-ready story — the everyday happy path through FastPPM.</div>"
        + _journey([
            ("Sign in", "Land on the chat analyst that already knows every "
                        "initiative, milestone, risk and value driver."),
            ("Ingest", "Drag a status report, tracker or deck (PDF / XLSX / PPTX / "
                       "DOCX) into <b>Documents</b> — the AI extracts the structure."),
            ("Review &amp; merge", "Check the <b>before → after</b>, then <b>merge "
                                    "into the master repository</b> in one click."),
            ("Track", "Your initiative is now live in the <b>registry, Gantt and "
                      "dashboard</b> — one source of truth."),
            ("Ask &amp; update", "Ask the copilot for status, risk and value — and "
                                  "give updates in plain English, with an audit trail."),
            ("Report", "Generate a <b>board-ready report</b> from a template, edit "
                       "it block-by-block, and export to PDF / DOCX / PPTX."),
        ])),
    ("The layout", _cards([
        ("Left", "Navigation, recent chats, and the initiative tree (by workstream). "
                 "Click any initiative to open it."),
        ("Centre", "The chat analyst (home), or whichever page you've opened."),
        ("Right", "The copilot rail (“Ask FastPPM”) and a live programme snapshot."),
    ])),
    ("1 · Ask FastPPM", _bullets([
        "The home screen is a chat analyst that knows all your data.",
        "<i>“What's the overall status of the Value Creation Plan?”</i>",
        "<i>“Which initiatives are red or at risk, and why?”</i>",
        "<i>“Show me the top value drivers by realised benefit.”</i>",
        "Give updates in plain English — applied with an audit entry:",
        "<i>“Mark Design / baseline as 80% complete.”</i> · <i>“Set ERP modernisation to at risk.”</i>",
    ])),
    ("2 · Ingest documents (flagship)", """
      <div class='flow'>
        <span>Upload</span><span>▶</span><span>Extract</span><span>▶</span>
        <span>Review</span><span>▶</span><span>Merge</span>
      </div>""" + _bullets([
        "Drag &amp; drop <b>PDF / XLSX / PPTX / DOCX</b> status reports, trackers and decks.",
        "The AI extracts <b>milestones, risks and value drivers</b>, normalising varied "
        "column names and date formats, and flags <b>inconsistencies</b>.",
        "<b>Review</b> the extracted entities, then <b>Merge into master repository</b> —",
        "it appears as an initiative in the registry, Gantt and dashboard, searchable from chat.",
    ])),
    ("Ask FastPPM — in action", _shot(
        "chat.png", "The chat analyst answers status, value and risk questions over your "
        "live data — and applies natural-language updates.")),
    ("Charts on demand", """
      <div class='flow'>
        <span>Ask for stats</span><span>▶</span><span>Table</span><span>▶</span>
        <span>Pick a chart</span><span>▶</span><span>Chart</span>
      </div>""" + _bullets([
        "Ask a quantitative question — <i>“value realised by category?”</i>, "
        "<i>“how many initiatives are red?”</i>, <i>“show the value trend”</i>.",
        "FastPPM answers with a <b>table of the real numbers</b>, then offers "
        "<b>chart chips</b> — Bar, Pie or Line — with the best one highlighted.",
        "Click a chip and the <b>interactive chart draws inline</b>; switch type "
        "any time. Time series default to a <b>line</b>, comparisons to a <b>bar</b>, "
        "parts-of-whole to a <b>pie</b>.",
    ])),
    ("Conversational analytics", _shot(
        "analytics.png", "Ask for statistics and FastPPM shows the table, then lets "
        "you turn it into a bar, pie or line chart in one click — straight in the chat.")),
    ("Documents — the ingestion queue", _shot(
        "documents.png", "Upload PDF / XLSX / PPTX / DOCX; each is auto-extracted with "
        "milestone / risk / value-driver counts and inconsistency flags.")),
    ("The transformation, made visible", _shot(
        "before-after.png", "Before → after on every document: your messy source on the "
        "left, the clean structured result on the right — then “Merge into master "
        "repository”.")),
    ("Report builder", _shot(
        "report-builder.png", "① Upload a report template, ② Generate a board-ready report "
        "that merges the programme data into it — with a live preview on the right.")),
    ("Edit &amp; export the report", _shot(
        "report-editor.png", "Edit the AI-assembled report block-by-block in a WYSIWYG "
        "editor (live preview), then export to PDF, DOCX or PPTX.")),
    ("3 · Track initiatives", _bullets([
        "<b>Initiatives</b> — every initiative &amp; AI use case, filterable by type and status.",
        "Open one for its <b>milestones</b> (baseline vs planned vs actual), <b>risk register</b>,",
        "<b>value</b>, the <b>source document</b> it was ingested from, and a full <b>audit trail</b>.",
    ])),
    ("Initiative detail", _shot(
        "initiative.png", "Progress, value, milestones, risks and audit trail for one "
        "initiative — and download any table as ⬇ CSV or ⬇ XLSX.")),
    ("4 · The master Gantt", _shot(
        "gantt.png", "RAG-coloured bars over their grey baseline · milestone diamonds by "
        "status · today line · Kanban &amp; dependency map below.")),
    ("5 · Programme dashboard", _shot(
        "dashboard.png", "On-track %, on-time delivery, value realised vs target, the value "
        "waterfall, RAG health and recent activity.")),
    ("Prompt Manager", _shot(
        "prompts.png", "Edit the extraction prompt in a WYSIWYG editor (no markdown "
        "needed) — every save is a new version you can re-activate.")),
    ("Sign in &amp; help", _bullets([
        "<b>Sign in</b> with Google or email &amp; password (authorised accounts only).",
        "<b>Help</b> → this User Guide (PDF / PPTX) and the Technical Architecture.",
        "Everything is a single source of truth — chat, Gantt, value and the master registry "
        "all read the same data.",
    ])),
]

# ── HTML deck ───────────────────────────────────────────────────────────────

CSS = """
*{margin:0;padding:0;box-sizing:border-box}
body{font-family:-apple-system,Segoe UI,Roboto,Helvetica,Arial,sans-serif;color:#48484f}
.slide{width:1280px;height:720px;position:relative;overflow:hidden;background:#fbfcfd;
  border-bottom:1px solid #fff}
.slide .hd{height:96px;background:linear-gradient(90deg,#123B5D,#00A6A6);color:#fff;
  display:flex;align-items:center;padding:0 56px;font-size:34px;font-weight:700}
.slide .bd{padding:44px 56px;font-size:23px;line-height:1.5}
.slide .bd ul{list-style:none}
.slide .bd li{margin:14px 0;padding-left:30px;position:relative}
.slide .bd li:before{content:"";position:absolute;left:0;top:12px;width:12px;height:12px;
  border-radius:3px;background:#00A6A6}
.slide .bd i{color:#123B5D;font-style:italic}
.foot-bar{position:absolute;bottom:0;left:0;right:0;height:46px;background:#0B2942;color:#D6E9EF;
  display:flex;align-items:center;justify-content:space-between;padding:0 56px;font-size:15px}
.cards{display:flex;gap:22px;margin-top:8px}
.card{flex:1;background:#fff;border:1px solid #D8E5EA;border-radius:14px;padding:24px}
.card .ct{font-weight:700;color:#123B5D;font-size:24px;margin-bottom:10px}
.card .cb{font-size:20px;color:#48484f;line-height:1.45}
.lead{font-size:21px;color:#123B5D;margin-bottom:22px;line-height:1.45}
.journey{display:grid;grid-template-columns:repeat(3,1fr);gap:20px}
.step{position:relative;background:#fff;border:1px solid #D8E5EA;border-radius:14px;
  padding:22px 22px 20px;box-shadow:0 3px 12px rgba(18,59,93,.06)}
.step .sn{position:absolute;top:-16px;left:22px;width:38px;height:38px;border-radius:10px;
  background:linear-gradient(135deg,#123B5D,#00A6A6);color:#fff;font-weight:800;font-size:20px;
  display:flex;align-items:center;justify-content:center;box-shadow:0 3px 8px rgba(18,59,93,.25)}
.step .st{font-weight:700;color:#123B5D;font-size:21px;margin:8px 0 8px}
.step .sb{font-size:17px;color:#48484f;line-height:1.42}
.step b{color:#00A6A6;font-weight:600}
.flow{display:flex;gap:14px;align-items:center;font-size:22px;font-weight:600;color:#123B5D;margin-bottom:18px}
.flow span:nth-child(even){color:#9AAEB8;font-weight:400}
.hero{height:100%;background:linear-gradient(135deg,#123B5D,#0B2942);color:#fff;
  display:flex;flex-direction:column;justify-content:center;padding:0 90px}
.hero .brand{font-size:92px;font-weight:800;letter-spacing:-1px}
.hero .brand span{color:#00A6A6}
.hero .sub{font-size:30px;color:#C4E3E8;margin-top:6px}
.hero .tag{font-size:23px;color:#DCEDEF;margin-top:34px;max-width:880px;line-height:1.5}
.hero .foot{position:absolute;bottom:54px;font-size:20px;color:#9FC8D2;letter-spacing:3px;text-transform:uppercase}
.slide .bd:has(.shot){padding:18px 56px 0}
.shot{height:486px;display:flex;align-items:center;justify-content:center}
.shot img{max-width:100%;max-height:486px;border:1px solid #D8E5EA;border-radius:10px;
  box-shadow:0 6px 22px rgba(18,59,93,.14)}
.cap{text-align:center;font-size:18px;color:#7a7a85;margin-top:10px;line-height:1.4}
"""


def _slide_html(idx, kind, body):
    if kind == "title":
        return f"<section class='slide'>{body}</section>"
    return (f"<section class='slide'><div class='hd'>{kind}</div>"
            f"<div class='bd'>{body}</div>"
            f"<div class='foot-bar'><span>FastPPM · User Guide</span>"
            f"<span>{idx} / {len(SLIDES)-1}</span></div></section>")


def build_html():
    body = "".join(_slide_html(i, k, b) for i, (k, b) in enumerate(SLIDES))
    return f"<!doctype html><html><head><meta charset='utf-8'><style>{CSS}</style></head>" \
           f"<body>{body}</body></html>"


# ── Render → PNG → PDF + PPTX ───────────────────────────────────────────────

def _screenshot_slides(html: str) -> list[bytes]:
    from playwright.sync_api import sync_playwright
    pngs = []
    with tempfile.TemporaryDirectory() as td:
        f = Path(td) / "deck.html"
        f.write_text(html)
        with sync_playwright() as p:
            b = p.chromium.launch()
            pg = b.new_page(viewport={"width": W, "height": H}, device_scale_factor=2)
            pg.goto(f.as_uri())
            for el in pg.query_selector_all(".slide"):
                pngs.append(el.screenshot())
            b.close()
    return pngs


def _write_pdf(pngs, out: Path):
    from fpdf import FPDF
    # 16:9 landscape in mm (1280:720 @ 96dpi). The format tuple (width > height)
    # already defines a landscape page; don't also pass orientation='L' (it would
    # swap the dimensions back to portrait).
    pw, ph = 338.667, 190.5
    pdf = FPDF(unit="mm", format=(pw, ph))
    pdf.set_auto_page_break(False)
    with tempfile.TemporaryDirectory() as td:
        for i, png in enumerate(pngs):
            p = Path(td) / f"s{i}.png"
            p.write_bytes(png)
            pdf.add_page()
            pdf.image(str(p), 0, 0, pw, ph)
        pdf.output(str(out))


def _write_pptx(pngs, out: Path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    with tempfile.TemporaryDirectory() as td:
        for i, png in enumerate(pngs):
            p = Path(td) / f"s{i}.png"
            p.write_bytes(png)
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(str(out))


BASENAME = "fastppm_user_guide"


def main():
    from datetime import date
    DOCS.mkdir(exist_ok=True)
    # RULE: every generated user guide is date-stamped in the filename so versions
    # are distinguishable (e.g. fastppm_user_guide_2026-06-22.pdf). The app serves
    # the newest dated file. Old dated outputs are pruned, keeping the latest only.
    stamp = date.today().isoformat()
    base = f"{BASENAME}_{stamp}"
    for old in DOCS.glob(f"{BASENAME}_*.*"):
        old.unlink()
    # Also drop any legacy non-dated artifacts.
    for legacy in ("html", "pdf", "pptx"):
        p = DOCS / f"{BASENAME}.{legacy}"
        if p.exists():
            p.unlink()

    html = build_html()
    (DOCS / f"{base}.html").write_text(html)
    # A date-stamped copy of the prose guide so the distributable trio
    # (md / pdf / pptx) is versioned together; the undated guide stays the in-app source.
    src_md = DOCS / "fastppm_user_guide.md"
    if src_md.exists():
        (DOCS / f"{base}.md").write_text(src_md.read_text())
    pngs = _screenshot_slides(html)
    _write_pdf(pngs, DOCS / f"{base}.pdf")
    _write_pptx(pngs, DOCS / f"{base}.pptx")
    print(f"Wrote {len(pngs)} slides → docs/{base}.{{md,html,pdf,pptx}}")


if __name__ == "__main__":
    main()
